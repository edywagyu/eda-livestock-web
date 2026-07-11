#!/usr/bin/env python3
"""EDA-LIVESTOCK Sales Deck Generator — matches website design system. V2 Complete Rewrite."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FOREST = RGBColor(0x0F, 0x3D, 0x2E)
FOREST_DEEP = RGBColor(0x0A, 0x2D, 0x21)
GOLD = RGBColor(0xD4, 0xA9, 0x3B)
GOLD_SOFT = RGBColor(0xB8, 0x93, 0x2F)
CREAM = RGBColor(0xFA, 0xF7, 0xF0)
PAPER = RGBColor(0xF5, 0xEF, 0xE3)
BONE = RGBColor(0xEF, 0xE8, 0xD7)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_MUTED = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_DISPLAY = "Arial Black"
FONT_BODY = "Arial"
FONT_JP = "Yu Gothic"
FONT_JP_SERIF = "Yu Mincho"

SLD_W = Inches(13.333)
SLD_H = Inches(7.5)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_name=FONT_BODY,
             font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, spacing_before=0, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.space_before = Pt(spacing_before)
    p.line_spacing = line_spacing
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox

def add_rich_text(slide, left, top, width, height, segments, alignment=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.line_spacing = line_spacing
    for i, seg in enumerate(segments):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            if seg.get("newline"):
                p = tf.add_paragraph()
                p.alignment = alignment
                p.line_spacing = line_spacing
                if seg.get("spacing_before"):
                    p.space_before = Pt(seg["spacing_before"])
                run = p.add_run()
            else:
                run = p.add_run()
        run.text = seg.get("text", "")
        run.font.name = seg.get("font", FONT_BODY)
        run.font.size = Pt(seg.get("size", 14))
        run.font.color.rgb = seg.get("color", WHITE)
        run.font.bold = seg.get("bold", False)
    return txBox

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    return shape

def add_line_shape(slide, left, top, width, color=GOLD, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_footer(slide, num, dark=True):
    c = RGBColor(0xFF, 0xFF, 0xFF) if dark else FOREST
    opacity_color = RGBColor(0x80, 0x80, 0x80) if dark else RGBColor(0x99, 0x99, 0x99)
    add_text(slide, Inches(0.6), Inches(7.05), Inches(3), Inches(0.35),
             "EDA-LIVESTOCK", FONT_DISPLAY, 9, opacity_color, alignment=PP_ALIGN.LEFT)
    add_text(slide, Inches(9.7), Inches(7.05), Inches(3), Inches(0.35),
             f"{num:02d}", FONT_DISPLAY, 9, opacity_color, alignment=PP_ALIGN.RIGHT)

def main():
    prs = Presentation()
    prs.slide_width = SLD_W
    prs.slide_height = SLD_H
    blank = prs.slide_layouts[6]

    # ════════════════════════════════════════════
    # SLIDE 01 — COVER
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST)
    add_text(s, Inches(0.8), Inches(1.2), Inches(8), Inches(3),
             "LIVESTOCK FOR\nTHE NEXT CENTURY.", FONT_DISPLAY, 60, WHITE, True)
    add_text(s, Inches(0.8), Inches(3.8), Inches(7), Inches(1),
             "無投薬 × 循環型飼料 × 世界初オーガニック認証和牛\n宮崎から世界13カ国・4大陸へ ―― 20代のチームが100年続く畜産を創る。",
             FONT_JP, 13, RGBColor(0xC0, 0xC0, 0xC0), line_spacing=1.8)
    add_line_shape(s, Inches(0.8), Inches(5.3), Inches(11.5), RGBColor(0x3A, 0x6A, 0x50))
    stats = [("13", "Countries · 4 Continents"), ("2023", "Founded in Miyazaki"),
             ("2026", "World's First Organic Wagyu"), ("1%", "Additive-Free Farms")]
    for i, (num, label) in enumerate(stats):
        x = Inches(0.8 + i * 3.0)
        add_text(s, x, Inches(5.5), Inches(2.8), Inches(0.7), num, FONT_DISPLAY, 40, GOLD, True)
        add_text(s, x, Inches(6.15), Inches(2.8), Inches(0.4), label, FONT_DISPLAY, 9,
                 RGBColor(0x80, 0x99, 0x88), alignment=PP_ALIGN.LEFT)
    add_text(s, Inches(7.5), Inches(6.8), Inches(5.5), Inches(0.3),
             "CONFIDENTIAL — EDA LIVESTOCK CO., LTD.", FONT_DISPLAY, 8,
             RGBColor(0x40, 0x60, 0x50), alignment=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════════
    # SLIDE 02 — VISION
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_rect(s, Inches(6.666), Inches(0), Inches(6.666), SLD_H, FOREST)
    add_text(s, Inches(0.6), Inches(0.8), Inches(5.5), Inches(0.3),
             "— OUR VISION —", FONT_DISPLAY, 11, GOLD_SOFT)
    add_rich_text(s, Inches(0.6), Inches(1.3), Inches(5.5), Inches(1.6), [
        {"text": "牛・生産者・消費者、\n", "font": FONT_JP, "size": 32, "color": FOREST, "bold": True},
        {"text": "三方よし", "font": FONT_JP, "size": 32, "color": GOLD, "bold": True},
        {"text": "を、和牛で。", "font": FONT_JP, "size": 32, "color": FOREST, "bold": True},
    ])
    add_text(s, Inches(0.6), Inches(3.2), Inches(5.5), Inches(2),
             "私たちは、和牛の生産を「三方よし」で再設計します。牛が穏やかに育ち、生産者が正当な対価を得て、消費者が安心して味わえる。この三者すべてが満足する仕組みこそ、100年続く畜産の条件です。\n\n既存のスタンダードを否定するのではなく、新しい選択肢を増やす。それが私たちの立ち位置です。",
             FONT_JP, 12, INK_MUTED, line_spacing=1.85)
    add_text(s, Inches(7.0), Inches(0.8), Inches(5.5), Inches(0.3),
             "— INDUSTRY CHALLENGE —", FONT_DISPLAY, 11, GOLD)
    add_text(s, Inches(7.0), Inches(1.3), Inches(3), Inches(1), "70", FONT_DISPLAY, 64, GOLD, True)
    add_text(s, Inches(9.0), Inches(1.55), Inches(2), Inches(0.5), "歳超", FONT_JP, 20,
             RGBColor(0x80, 0x80, 0x80))
    add_text(s, Inches(7.0), Inches(2.3), Inches(5), Inches(0.3),
             "日本の畜産業界 平均年齢", FONT_JP, 11, RGBColor(0x80, 0x99, 0x88))
    cards = [
        ("01 — CATTLE", "For the Cattle", "化学物質に頼らず、十分な空間と時間で穏やかに育つ環境。"),
        ("02 — PRODUCER", "For the Producer", "直接輸出と適正価格で、生産者がきちんと利益を得られる仕組み。"),
        ("03 — CONSUMER", "For the Consumer", "安全で味わい深く、透明性のある肉を世界の食卓に。"),
    ]
    for i, (num, title, desc) in enumerate(cards):
        y = Inches(3.1 + i * 1.35)
        add_line_shape(s, Inches(7.0), y, Pt(2), GOLD, Pt(2))
        add_rect(s, Inches(7.0), y, Pt(2), Inches(1.1), GOLD)
        add_text(s, Inches(7.2), y, Inches(5), Inches(0.25), num, FONT_DISPLAY, 9, GOLD)
        add_text(s, Inches(7.2), Inches(y.inches + 0.3), Inches(5), Inches(0.3),
                 title, FONT_DISPLAY, 16, WHITE, True)
        add_text(s, Inches(7.2), Inches(y.inches + 0.6), Inches(5), Inches(0.5),
                 desc, FONT_JP, 10, RGBColor(0xBB, 0xBB, 0xBB), line_spacing=1.7)
    add_footer(s, 2, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 03 — PHILOSOPHY / HOW WE RAISED
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST_DEEP)
    add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.3),
             "— OUR PHILOSOPHY —", FONT_DISPLAY, 11, GOLD, alignment=PP_ALIGN.CENTER)
    add_rich_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8), [
        {"text": "NOT THE RANK, ", "font": FONT_DISPLAY, "size": 36, "color": WHITE, "bold": True},
        {"text": "BUT HOW THEY'RE RAISED.", "font": FONT_DISPLAY, "size": 36, "color": GOLD, "bold": True},
    ], alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(2.5), Inches(1.9), Inches(8.3), Inches(0.6),
             "A5・BMS・産地名 ーー これまで和牛は数字とラベルで語られてきました。私たちが大切にするのは、その手前にある「どう育てたか」です。",
             FONT_JP, 11, RGBColor(0x88, 0x99, 0x90), alignment=PP_ALIGN.CENTER, line_spacing=1.7)
    pillars = [
        ("01 / CIRCULAR", "CIRCULAR FEED", "循環型飼料",
         "自社の堆肥から飼料を育て、再び牛に与える。輸入飼料に頼らない、国産循環型の畜産を実践。"),
        ("02 / PURE", "NO CHEMICAL", "抗生物質 · ホルモン剤 ゼロ",
         "抗生物質・成長促進剤・ホルモン剤いずれも一切不使用。日本でこれを徹底する農家はごくわずか。"),
        ("03 / WELFARE", "ANIMAL WELFARE", "アニマルウェルフェア",
         "有機規格に準じた広い牛舎。ストレスの少ない環境でのびのびと育成。"),
        ("04 / NATURAL", "MINERAL WATER", "霧島連山の天然水",
         "霧島連山の伏流水を全頭の飲み水に使用。ミネラルが肉の旨みを底上げ。"),
    ]
    for i, (num, title, jp, desc) in enumerate(pillars):
        x = Inches(0.6 + i * 3.1)
        add_rect(s, x, Inches(2.9), Inches(2.9), Inches(4.0),
                 RGBColor(0x0E, 0x38, 0x2A), GOLD, Pt(0.5))
        add_line_shape(s, x, Inches(2.9), Inches(2.9), GOLD, Pt(2))
        add_text(s, Inches(x.inches + 0.2), Inches(3.15), Inches(2.5), Inches(0.3),
                 num, FONT_DISPLAY, 10, GOLD)
        add_text(s, Inches(x.inches + 0.2), Inches(3.8), Inches(2.5), Inches(0.4),
                 title, FONT_DISPLAY, 18, WHITE, True)
        add_text(s, Inches(x.inches + 0.2), Inches(4.3), Inches(2.5), Inches(0.3),
                 jp, FONT_JP, 9, GOLD_SOFT)
        add_text(s, Inches(x.inches + 0.2), Inches(4.8), Inches(2.5), Inches(1.5),
                 desc, FONT_JP, 10, RGBColor(0xBB, 0xBB, 0xBB), line_spacing=1.75)
    add_footer(s, 3)

    # ════════════════════════════════════════════
    # SLIDE 04 — BRAND ARCHITECTURE
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.3),
             "— BRAND ARCHITECTURE —", FONT_DISPLAY, 11, GOLD_SOFT)
    add_rich_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.8), [
        {"text": "飼料の規定でブランド化する、", "font": FONT_JP, "size": 26, "color": FOREST, "bold": True},
        {"text": "日本唯一の和牛。", "font": FONT_JP, "size": 26, "color": GOLD, "bold": True},
    ])
    add_text(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.3),
             "産地ではなく「飼料」と「飼育方法」でブランドを規定。等級指定可能、産地指定は金額UP。",
             FONT_JP, 11, INK_MUTED)
    eda_specs = [
        ("FEED", "循環型飼料使用（自社堆肥起点）"),
        ("CHEMICAL", "抗生物質 · ホルモン剤 不使用"),
        ("GRADE", "A3以上 · 28ヶ月以上肥育"),
        ("REGION", "宮崎 · 鹿児島 · 東北（北海道 · 青森 · 秋田）"),
        ("WATER", "天然水（霧島連山伏流水）"),
        ("OPTION", "等級指定可 / 産地指定は金額UP"),
    ]
    add_rect(s, Inches(0.6), Inches(2.15), Inches(5.8), Inches(4.0), None, BONE, Pt(1))
    add_text(s, Inches(0.9), Inches(2.3), Inches(5), Inches(0.5),
             "Eda Wagyu", FONT_DISPLAY, 24, FOREST, True)
    add_text(s, Inches(0.9), Inches(2.75), Inches(5), Inches(0.3),
             "江田和牛", FONT_JP, 11, GOLD_SOFT)
    for i, (label, val) in enumerate(eda_specs):
        y = Inches(3.2 + i * 0.48)
        add_text(s, Inches(0.9), y, Inches(1.2), Inches(0.35), label, FONT_DISPLAY, 9, GOLD_SOFT)
        add_text(s, Inches(2.2), y, Inches(4), Inches(0.35), val, FONT_JP, 11, INK)

    prem_specs = [
        ("ORGANIC", "有機JAS認証 · EU Leaf 同等性承認"),
        ("FEED", "100% 有機認証飼料"),
        ("STATUS", "世界初のオーガニック認証黒毛和牛"),
        ("LAUNCH", "2026年2月 ローンチ"),
        ("SALES", "江田和牛お客様限定販売"),
        ("ONLY", "江田畜産のみが生産・販売"),
    ]
    add_rect(s, Inches(6.9), Inches(2.15), Inches(5.8), Inches(4.0), None, GOLD, Pt(1))
    add_rect(s, Inches(10.5), Inches(2.0), Inches(2.0), Inches(0.3), GOLD)
    add_text(s, Inches(10.5), Inches(2.0), Inches(2.0), Inches(0.3),
             "WORLD'S FIRST", FONT_DISPLAY, 9, FOREST, True, PP_ALIGN.CENTER)
    add_text(s, Inches(7.2), Inches(2.3), Inches(5), Inches(0.5),
             "Premium EdaWagyu", FONT_DISPLAY, 24, FOREST, True)
    add_text(s, Inches(7.2), Inches(2.75), Inches(5), Inches(0.3),
             "特選江田和牛 · 有機JAS認証", FONT_JP, 11, GOLD_SOFT)
    for i, (label, val) in enumerate(prem_specs):
        y = Inches(3.2 + i * 0.48)
        add_text(s, Inches(7.2), y, Inches(1.2), Inches(0.35), label, FONT_DISPLAY, 9, GOLD_SOFT)
        add_text(s, Inches(8.5), y, Inches(4), Inches(0.35), val, FONT_JP, 11, INK)

    add_rect(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.6), FOREST)
    add_text(s, Inches(0.8), Inches(6.42), Inches(11.7), Inches(0.55),
             "日本の和牛ブランドで「飼料の規定」によるブランド化を実現しているのは江田畜産のみ。オーガニック和牛は世界で唯一、江田畜産だけが生産。",
             FONT_JP, 11, WHITE, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 4, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 05 — COMPANY PROFILE
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_rect(s, Inches(7.0), Inches(0), Inches(6.333), SLD_H, FOREST)
    add_text(s, Inches(0.6), Inches(0.5), Inches(6), Inches(0.3),
             "— COMPANY PROFILE —", FONT_DISPLAY, 11, GOLD_SOFT)
    add_text(s, Inches(0.6), Inches(0.95), Inches(6), Inches(0.5),
             "会社概要", FONT_JP, 26, FOREST, True)
    profile_rows = [
        ("COMPANY", "江田畜産株式会社 / Eda Livestock Co., Ltd."),
        ("REPRESENTATIVE", "代表取締役 江田 友輝（Tomoki Eda）"),
        ("FOUNDED", "2023年4月"),
        ("BUSINESS", "黒毛和牛 繁殖 · 肥育 · 小売 · 輸出 · スマート畜産"),
        ("FARMS", "宮崎 · 鹿児島 · 東北（北海道 · 青森 · 秋田）"),
        ("HQ", "〒889-4411 宮崎県西諸県郡高原町大字広原 2025-2"),
        ("BRANDS", "Eda Wagyu / Premium EdaWagyu（有機JAS · EU Leaf同等性）"),
        ("EXPORT", "世界13カ国 · 4大陸 · 直接輸出100%"),
        ("SETTLEMENT", "円建て（JPY）"),
        ("PARTNERS", "SoftBank · TOYOTA (TTDC) · 内閣府クールジャパン"),
    ]
    for i, (label, val) in enumerate(profile_rows):
        y = Inches(1.7 + i * 0.5)
        add_line_shape(s, Inches(0.6), Inches(y.inches + 0.42), Inches(6.0),
                       RGBColor(0xE8, 0xE0, 0xD0), Pt(0.5))
        add_text(s, Inches(0.6), y, Inches(1.5), Inches(0.4), label, FONT_DISPLAY, 9, GOLD_SOFT)
        add_text(s, Inches(2.2), y, Inches(4.4), Inches(0.4), val, FONT_JP, 11, INK)

    add_text(s, Inches(7.3), Inches(0.6), Inches(5.5), Inches(0.3),
             "— INTEGRATED FARMING —", FONT_DISPLAY, 11, GOLD)
    add_text(s, Inches(7.3), Inches(1.05), Inches(5.5), Inches(0.5),
             "一貫農業フロー", FONT_JP, 22, WHITE, True)
    flow = [
        ("01", "分娩", "Birth"),
        ("02", "繁殖（10ヶ月前後）", "Breeding"),
        ("03", "肥育（20ヶ月以上）", "Fattening · 28ヶ月+"),
        ("04", "出荷 → 世界13カ国へ直接輸出", "Direct Export · 円建て決済"),
    ]
    for i, (num, jp, en) in enumerate(flow):
        y = Inches(1.9 + i * 1.15)
        add_rect(s, Inches(7.3), y, Inches(0.45), Inches(0.45), None, GOLD, Pt(1))
        add_text(s, Inches(7.3), y, Inches(0.45), Inches(0.45), num,
                 FONT_DISPLAY, 14, GOLD, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.0), y, Inches(4.5), Inches(0.3), jp, FONT_JP, 12, WHITE)
        add_text(s, Inches(8.0), Inches(y.inches + 0.3), Inches(4.5), Inches(0.25),
                 en, FONT_DISPLAY, 9, RGBColor(0x80, 0x99, 0x88))
        if i < 3:
            add_text(s, Inches(7.4), Inches(y.inches + 0.55), Inches(0.3), Inches(0.3),
                     "↓", FONT_BODY, 14, RGBColor(0x50, 0x70, 0x60), alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(7.3), Inches(6.2), Inches(5.5), Inches(0.8),
             "繁殖と肥育の両方を経営する一貫農業。仔牛の段階から飼料と環境を管理し、トレーサビリティを完全担保。",
             FONT_JP, 10, RGBColor(0x88, 0x99, 0x90), line_spacing=1.7)
    add_footer(s, 5, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 06 — ADOPTED BY
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST)
    add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.3),
             "— WHERE EDA WAGYU IS SERVED —", FONT_DISPLAY, 11, GOLD, alignment=PP_ALIGN.CENTER)
    add_rich_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8), [
        {"text": "TRUSTED BY ", "font": FONT_DISPLAY, "size": 36, "color": WHITE, "bold": True},
        {"text": "THE WORLD'S FINEST.", "font": FONT_DISPLAY, "size": 36, "color": GOLD, "bold": True},
    ], alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(2.5), Inches(1.85), Inches(8.3), Inches(0.3),
             "日本国内外の三ツ星ホテル・高級百貨店で採用。創業3年で21拠点に拡大。",
             FONT_JP, 11, RGBColor(0x80, 0x99, 0x88), alignment=PP_ALIGN.CENTER)
    hotels = [
        ("The Ritz-Carlton", "Tokyo · Fukuoka · HK · NY", True),
        ("Four Seasons Hotel", "Tokyo at Otemachi", True),
        ("The Peninsula", "Tokyo", True),
        ("Mandarin Oriental", "Tokyo", True),
        ("Tokyo Station Hotel", "日本重要文化財 · 東京駅", False),
        ("Tokyo Edition Ginza", "Marriott Luxury · 銀座", False),
        ("city'super", "Hong Kong 全店展開", False),
        ("Mitsukoshi · Keio", "百貨店 · 銀座 · 日本橋 · 新宿", False),
    ]
    for i, (name, loc, featured) in enumerate(hotels):
        col = i % 4
        row = i // 4
        x = Inches(0.6 + col * 3.1)
        y = Inches(2.6 + row * 2.2)
        border = GOLD if featured else RGBColor(0x20, 0x50, 0x3C)
        bg = RGBColor(0x14, 0x48, 0x35) if featured else RGBColor(0x12, 0x42, 0x30)
        add_rect(s, x, y, Inches(2.9), Inches(1.9), bg, border, Pt(1))
        add_text(s, Inches(x.inches + 0.2), Inches(y.inches + 0.5), Inches(2.5), Inches(0.7),
                 name, FONT_DISPLAY, 16, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(x.inches + 0.2), Inches(y.inches + 1.3), Inches(2.5), Inches(0.4),
                 loc, FONT_JP, 9, RGBColor(0x80, 0x99, 0x88), alignment=PP_ALIGN.CENTER)
    add_footer(s, 6)

    # ════════════════════════════════════════════
    # SLIDE 07 — GLOBAL PRESENCE
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST_DEEP)
    add_text(s, Inches(0.6), Inches(0.5), Inches(6), Inches(0.3),
             "— GLOBAL PRESENCE —", FONT_DISPLAY, 11, GOLD)
    add_rich_text(s, Inches(0.6), Inches(0.95), Inches(6), Inches(0.8), [
        {"text": "FARM TO ", "font": FONT_DISPLAY, "size": 36, "color": WHITE, "bold": True},
        {"text": "THE WORLD.", "font": FONT_DISPLAY, "size": 36, "color": GOLD, "bold": True},
    ])
    add_text(s, Inches(7.5), Inches(0.7), Inches(5.2), Inches(0.6),
             "商社を介さない直接輸出100%。\n創業1年で1カ国、2年で7カ国、3年で13カ国。",
             FONT_JP, 10, RGBColor(0x80, 0x99, 0x88), alignment=PP_ALIGN.RIGHT, line_spacing=1.7)
    big_stats = [("13", "COUNTRIES"), ("4", "CONTINENTS"), ("100%", "DIRECT EXPORT")]
    for i, (num, label) in enumerate(big_stats):
        x = Inches(0.6 + i * 2.8)
        add_text(s, x, Inches(2.0), Inches(2.5), Inches(0.8), num, FONT_DISPLAY, 48, GOLD, True)
        add_text(s, x, Inches(2.75), Inches(2.5), Inches(0.3), label, FONT_DISPLAY, 10,
                 RGBColor(0x80, 0x99, 0x88))
    regions = [
        ("ASIA · 7 COUNTRIES", "🇭🇰 Hong Kong  🇸🇬 Singapore  🇹🇭 Thailand  🇻🇳 Vietnam  🇹🇼 Taiwan  🇲🇾 Malaysia  🇮🇩 Indonesia"),
        ("EUROPE · 4 COUNTRIES", "🇩🇪 Germany  🇫🇮 Finland  🇪🇸 Spain  🇮🇹 Italy"),
        ("AMERICAS · 1 COUNTRY", "🇺🇸 USA (New York)"),
        ("OCEANIA · 1 COUNTRY", "🇦🇺 Australia"),
    ]
    for i, (label, countries) in enumerate(regions):
        y = Inches(3.5 + i * 0.9)
        add_text(s, Inches(0.6), y, Inches(4), Inches(0.25), label, FONT_DISPLAY, 11, GOLD)
        add_line_shape(s, Inches(0.6), Inches(y.inches + 0.3), Inches(6), RGBColor(0x18, 0x40, 0x30))
        add_text(s, Inches(0.6), Inches(y.inches + 0.4), Inches(12), Inches(0.35),
                 countries, FONT_JP, 10, RGBColor(0xBB, 0xBB, 0xBB))
    bars = [(1, "YEAR 1", 0.4), (7, "YEAR 2", 1.3), (13, "YEAR 3", 2.3)]
    for i, (num, yr, h) in enumerate(bars):
        x = Inches(8.5 + i * 1.6)
        add_rect(s, x, Inches(6.8 - h), Inches(0.8), Inches(h), GOLD)
        add_text(s, Inches(x.inches - 0.3), Inches(6.85), Inches(1.4), Inches(0.25),
                 yr, FONT_DISPLAY, 10, RGBColor(0x80, 0x99, 0x88), alignment=PP_ALIGN.CENTER)
        add_text(s, Inches(x.inches - 0.3), Inches(6.55 - h), Inches(1.4), Inches(0.3),
                 str(num), FONT_DISPLAY, 18, GOLD, True, PP_ALIGN.CENTER)
    add_footer(s, 7)

    # ════════════════════════════════════════════
    # SLIDE 08 — CERTIFICATIONS
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.3),
             "— CERTIFICATIONS & STANDARDS —", FONT_DISPLAY, 11, GOLD_SOFT)
    add_rich_text(s, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7), [
        {"text": "世界基準の認証群で、", "font": FONT_JP, "size": 24, "color": FOREST, "bold": True},
        {"text": "信頼をデータで証明する。", "font": FONT_JP, "size": 24, "color": GOLD, "bold": True},
    ])
    certs = [
        ("☪", "HALAL", "マレーシア · インドネシア\nイスラム圏への輸出対応"),
        ("🇺🇸", "USDA ORGANIC", "アメリカ農務省\n有機認証（申請中）"),
        ("🇪🇺", "EU LEAF 同等性", "欧州連合オーガニック\n同等性承認済み"),
        ("🍀", "有機JAS", "世界初 オーガニック認証\n黒毛和牛（2026年1月取得）"),
    ]
    for i, (icon, name, desc) in enumerate(certs):
        x = Inches(0.6 + i * 3.1)
        add_rect(s, x, Inches(1.9), Inches(2.9), Inches(1.7), WHITE, BONE, Pt(1))
        add_text(s, x, Inches(2.0), Inches(2.9), Inches(0.5), icon,
                 FONT_BODY, 24, FOREST, alignment=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.5), Inches(2.9), Inches(0.3), name,
                 FONT_DISPLAY, 14, FOREST, True, PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.85), Inches(2.9), Inches(0.6), desc,
                 FONT_JP, 9, INK_MUTED, alignment=PP_ALIGN.CENTER, line_spacing=1.6)

    compare = [
        ("項目", "EDA WAGYU", "PREMIUM EDA", "一般的な和牛"),
        ("抗生物質 · ホルモン剤", "✓ 不使用", "✓ 不使用", "使用あり"),
        ("飼料", "循環型飼料", "有機認証飼料100%", "輸入配合飼料"),
        ("オーガニック認証", "—", "✓ 世界唯一", "—"),
        ("アニマルウェルフェア", "✓ 対応", "✓ 有機規格準拠", "基準なし"),
        ("HALAL対応", "✓ 対応", "✓ 対応", "一部のみ"),
    ]
    col_widths = [Inches(3.0), Inches(3.0), Inches(3.0), Inches(3.0)]
    for row_i, row in enumerate(compare):
        y = Inches(3.85 + row_i * 0.48)
        if row_i == 0:
            add_rect(s, Inches(0.6), y, Inches(3.0), Inches(0.42), None)
            add_rect(s, Inches(3.6), y, Inches(3.0), Inches(0.42), FOREST)
            add_rect(s, Inches(6.6), y, Inches(3.0), Inches(0.42), FOREST)
            add_rect(s, Inches(9.6), y, Inches(3.0), Inches(0.42), BONE)
        for col_i, cell in enumerate(row):
            x = Inches(0.6 + col_i * 3.0)
            if row_i == 0:
                c = WHITE if col_i in (1, 2) else (INK_MUTED if col_i == 3 else INK)
                add_text(s, x, y, Inches(3.0), Inches(0.4), cell,
                         FONT_DISPLAY, 10, c, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            else:
                if row_i > 0:
                    add_line_shape(s, Inches(0.6), y, Inches(11.9),
                                  RGBColor(0xE8, 0xE0, 0xD0), Pt(0.3))
                c = FOREST if "✓" in cell else (RGBColor(0xCC, 0xCC, 0xCC) if cell in ("—", "使用あり", "輸入配合飼料", "基準なし", "一部のみ") else INK)
                al = PP_ALIGN.CENTER if col_i > 0 else PP_ALIGN.LEFT
                add_text(s, x, y, Inches(3.0), Inches(0.4), cell,
                         FONT_JP, 10, c, "✓" in cell, al, MSO_ANCHOR.MIDDLE)
    add_footer(s, 8, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 09 — BUSINESS MODEL
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST)
    add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.3),
             "— BUSINESS MODEL —", FONT_DISPLAY, 11, GOLD, alignment=PP_ALIGN.CENTER)
    add_rich_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7), [
        {"text": "DOMESTIC FEED → ", "font": FONT_DISPLAY, "size": 36, "color": WHITE, "bold": True},
        {"text": "GLOBAL EXPORT", "font": FONT_DISPLAY, "size": 36, "color": GOLD, "bold": True},
    ], alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(2.5), Inches(1.8), Inches(8.3), Inches(0.3),
             "国内で飼料をつくり、世界へ出荷する。円安が進むほど、このモデルは強くなる。",
             FONT_JP, 11, RGBColor(0x80, 0x99, 0x88), alignment=PP_ALIGN.CENTER)
    flow_items = [
        ("自社堆肥", "Compost", False),
        ("循環型飼料", "Circular Feed", False),
        ("一貫農業", "Integrated Farming", False),
        ("直接輸出", "13 Countries · 円建て", True),
        ("再投資", "Reinvest", True),
    ]
    for i, (jp, en, hl) in enumerate(flow_items):
        x = Inches(0.4 + i * 2.5)
        bg = RGBColor(0x1A, 0x50, 0x3A) if hl else RGBColor(0x12, 0x42, 0x30)
        bd = GOLD if hl else RGBColor(0x25, 0x55, 0x40)
        add_rect(s, x, Inches(2.6), Inches(2.1), Inches(1.2), bg, bd, Pt(1))
        add_text(s, x, Inches(2.7), Inches(2.1), Inches(0.5), jp,
                 FONT_JP, 14, WHITE, True, PP_ALIGN.CENTER)
        add_text(s, x, Inches(3.2), Inches(2.1), Inches(0.4), en,
                 FONT_DISPLAY, 9, RGBColor(0x80, 0x99, 0x88), alignment=PP_ALIGN.CENTER)
        if i < 4:
            add_text(s, Inches(x.inches + 2.1), Inches(2.9), Inches(0.4), Inches(0.5),
                     "→", FONT_DISPLAY, 20, RGBColor(0x50, 0x70, 0x60), alignment=PP_ALIGN.CENTER)
    points = [
        ("国内飼料生産", "循環農業のスケールアップ。北海道・九州の有機認証農家との連携。輸入依存からの脱却。"),
        ("円建て決済（JPY）", "全取引を円建てで決済。為替リスクを排除し、安定した収益構造を実現。外貨獲得型モデル。"),
        ("SMART AGRICULTURE", "SoftBank・TOYOTA(TTDC)と共同でAI画像解析・IoTセンシングによるスマート畜産を実証実験中。"),
    ]
    for i, (title, desc) in enumerate(points):
        x = Inches(0.6 + i * 4.1)
        add_line_shape(s, x, Inches(4.5), Inches(3.8), GOLD, Pt(2))
        add_text(s, x, Inches(4.7), Inches(3.8), Inches(0.3), title,
                 FONT_DISPLAY, 14, WHITE, True)
        add_text(s, x, Inches(5.15), Inches(3.8), Inches(1.5), desc,
                 FONT_JP, 10, RGBColor(0xBB, 0xBB, 0xBB), line_spacing=1.7)
    add_footer(s, 9)

    # ════════════════════════════════════════════
    # SLIDE 10 — TEAM
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    add_text(s, Inches(0.6), Inches(0.4), Inches(8), Inches(0.3),
             "— OUR TEAM —", FONT_DISPLAY, 11, GOLD_SOFT)
    add_rich_text(s, Inches(0.6), Inches(0.8), Inches(8), Inches(0.8), [
        {"text": "平均年齢70歳超の業界に、", "font": FONT_JP, "size": 22, "color": FOREST, "bold": True},
        {"text": "20代のチーム", "font": FONT_JP, "size": 22, "color": GOLD, "bold": True},
        {"text": "が挑む。", "font": FONT_JP, "size": 22, "color": FOREST, "bold": True},
    ])
    add_text(s, Inches(8.0), Inches(1.0), Inches(4.7), Inches(0.6),
             "戦略コンサル出身、海外大学卒、国際営業、有機JAS推進 ―― 多様なバックグラウンドが強み。",
             FONT_JP, 10, INK_MUTED, alignment=PP_ALIGN.RIGHT, line_spacing=1.6)
    members = [
        ("TE", "FOUNDER · CEO", "Tomoki Eda", "江田 友輝",
         "宮崎出身。日大商学部在学中にNZ・米留学。21歳で独立。27歳で江田畜産を創業、100年先の畜産を設計。", True),
        ("RT", "CO-FOUNDER · COO", "Rei Tanaka", "田中 伶",
         "早稲田大学卒。ベイカレント・コンサルティングで戦略・DXコンサルに従事。2023年共同創業。財務・オペレーション統括。", True),
        ("VQ", "INTERNATIONAL SALES", "Valjon Qejvani", "バルヨン · ケジヴァニ",
         "ウェスタンミシガン大学卒。米大手マーケティング会社を経て2023年ジョイン。HK・NY・EU開拓。", False),
        ("KN", "DOMESTIC · JAS LEAD", "Keigo Nono", "野々 圭吾",
         "専修大学卒。国内営業と有機JAS認証プロジェクトリーダー。世界初オーガニック和牛の実装を牽引。", False),
    ]
    for i, (init, role, name, name_jp, bio, founder) in enumerate(members):
        x = Inches(0.6 + i * 3.1)
        bd = GOLD if founder else BONE
        add_rect(s, x, Inches(1.9), Inches(2.9), Inches(4.0), WHITE, bd, Pt(1))
        add_rect(s, Inches(x.inches + 0.9), Inches(2.1), Inches(1.1), Inches(1.1), FOREST)
        add_text(s, Inches(x.inches + 0.9), Inches(2.1), Inches(1.1), Inches(1.1),
                 init, FONT_DISPLAY, 22, GOLD, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(x.inches + 0.2), Inches(3.4), Inches(2.5), Inches(0.25),
                 role, FONT_DISPLAY, 9, GOLD_SOFT, alignment=PP_ALIGN.CENTER)
        add_text(s, Inches(x.inches + 0.2), Inches(3.7), Inches(2.5), Inches(0.35),
                 name, FONT_DISPLAY, 16, FOREST, True, PP_ALIGN.CENTER)
        add_text(s, Inches(x.inches + 0.2), Inches(4.05), Inches(2.5), Inches(0.25),
                 name_jp, FONT_JP, 10, INK_MUTED, alignment=PP_ALIGN.CENTER)
        add_text(s, Inches(x.inches + 0.3), Inches(4.4), Inches(2.3), Inches(1.3),
                 bio, FONT_JP, 9, INK_MUTED, line_spacing=1.65)

    add_rect(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.8), WHITE, BONE, Pt(1))
    add_rect(s, Inches(0.8), Inches(6.25), Inches(0.7), Inches(0.6), FOREST)
    add_text(s, Inches(0.8), Inches(6.25), Inches(0.7), Inches(0.6),
             "HS", FONT_DISPLAY, 14, GOLD, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.7), Inches(6.2), Inches(4), Inches(0.25),
             "PROJECT DIRECTOR · CHEF ADVISOR", FONT_DISPLAY, 8, GOLD_SOFT)
    add_text(s, Inches(1.7), Inches(6.45), Inches(4), Inches(0.3),
             "Hitoshi Sugiura　杉浦 仁志", FONT_DISPLAY, 13, FOREST, True)
    add_text(s, Inches(6.5), Inches(6.25), Inches(6), Inches(0.6),
             "内閣府クールジャパンプロデューサー。\"Social Food Gastronomy\" を提唱。国賓レセプション代表料理人。",
             FONT_JP, 9, INK_MUTED, line_spacing=1.55)
    add_footer(s, 10, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 11 — JOURNEY
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST_DEEP)
    add_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.3),
             "— OUR JOURNEY —", FONT_DISPLAY, 11, GOLD, alignment=PP_ALIGN.CENTER)
    add_rich_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7), [
        {"text": "2023 — ", "font": FONT_DISPLAY, "size": 36, "color": WHITE, "bold": True},
        {"text": "2026", "font": FONT_DISPLAY, "size": 36, "color": GOLD, "bold": True},
    ], alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(3), Inches(1.55), Inches(7.3), Inches(0.3),
             "創業からわずか3年。世界が注目する和牛ブランドへ。",
             FONT_JP, 11, RGBColor(0x80, 0x99, 0x88), alignment=PP_ALIGN.CENTER)
    timeline = [
        ("2023", "YEAR 1 — FOUNDATION", [
            "江田畜産株式会社 創業（宮崎県高原町）",
            "Eda Wagyu ブランド確立",
            "香港 city'super 直接輸出開始",
            "無投薬・循環型飼料の体制構築",
        ], False),
        ("2024", "YEAR 2 — EXPANSION", [
            "輸出先 1→7カ国に拡大",
            "Ritz-Carlton · Four Seasons · Peninsula 採用",
            "東京ステーションホテル 採用",
            "SoftBank · TOYOTA(TTDC) スマート畜産実証",
        ], False),
        ("2025", "YEAR 3 — ACCELERATION", [
            "輸出先 7→13カ国に拡大",
            "EU圏（独·芬·西·伊）本格展開",
            "Malaysia · Indonesia HALAL対応",
            "Ritz-Carlton Hong Kong · New York 採用",
        ], False),
        ("2026", "NOW — WORLD'S FIRST", [
            "世界初 有機JAS認証 黒毛和牛 取得",
            "Premium EdaWagyu ローンチ",
            "EU Leaf 同等性 承認",
            "オーストラリア市場 進出",
        ], True),
    ]
    for i, (year, label, items, current) in enumerate(timeline):
        x = Inches(0.6 + i * 3.1)
        if current:
            add_rect(s, x, Inches(2.1), Inches(2.9), Inches(4.8),
                     RGBColor(0x12, 0x3A, 0x2A), GOLD, Pt(1))
        add_line_shape(s, x, Inches(2.2), Inches(2.9), GOLD, Pt(current and 3 or 2))
        add_text(s, Inches(x.inches + 0.2), Inches(2.4), Inches(2.5), Inches(0.6),
                 year, FONT_DISPLAY, 32, GOLD, True)
        add_text(s, Inches(x.inches + 0.2), Inches(3.0), Inches(2.5), Inches(0.3),
                 label, FONT_DISPLAY, 11, WHITE)
        for j, item in enumerate(items):
            c = GOLD if current and j < 2 else RGBColor(0xBB, 0xBB, 0xBB)
            add_text(s, Inches(x.inches + 0.2), Inches(3.6 + j * 0.7), Inches(2.5), Inches(0.6),
                     item, FONT_JP, 10, c, current and j < 2, line_spacing=1.5)
    add_footer(s, 11)

    # ════════════════════════════════════════════
    # SLIDE 12 — CONTACT
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST)
    add_rich_text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.0), [
        {"text": "LET'S BUILD ", "font": FONT_DISPLAY, "size": 42, "color": WHITE, "bold": True},
        {"text": "THE FUTURE TOGETHER.", "font": FONT_DISPLAY, "size": 42, "color": GOLD, "bold": True},
    ], alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(3), Inches(2.7), Inches(7.3), Inches(0.5),
             "お問い合わせ・ご商談・サンプルリクエストは、以下までお気軽にご連絡ください。",
             FONT_JP, 12, RGBColor(0x88, 0x99, 0x90), alignment=PP_ALIGN.CENTER)
    add_line_shape(s, Inches(3.5), Inches(3.5), Inches(2.8), RGBColor(0x30, 0x60, 0x48))
    add_text(s, Inches(3.5), Inches(3.7), Inches(2.8), Inches(0.3),
             "HEADQUARTERS", FONT_DISPLAY, 12, GOLD)
    add_text(s, Inches(3.5), Inches(4.1), Inches(2.8), Inches(2),
             "江田畜産株式会社\nEda Livestock Co., Ltd.\n\n〒889-4411\n宮崎県西諸県郡高原町大字広原 2025-2\n\nbackoffice@eda-livestock.com",
             FONT_JP, 11, RGBColor(0xBB, 0xBB, 0xBB), line_spacing=1.7)
    add_line_shape(s, Inches(7.0), Inches(3.5), Inches(3.0), RGBColor(0x30, 0x60, 0x48))
    add_text(s, Inches(7.0), Inches(3.7), Inches(3.0), Inches(0.3),
             "INTERNATIONAL SALES", FONT_DISPLAY, 12, GOLD)
    add_text(s, Inches(7.0), Inches(4.1), Inches(3.5), Inches(2),
             "Export Inquiry · Sample Request · Partnership\n\nbackoffice@eda-livestock.com\n\nInstagram: @eda_livestock\nWeb: eda-livestock.com",
             FONT_JP, 11, RGBColor(0xBB, 0xBB, 0xBB), line_spacing=1.7)
    add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.3),
             "© 2026 EDA-LIVESTOCK CO., LTD. — CONFIDENTIAL",
             FONT_DISPLAY, 9, RGBColor(0x30, 0x50, 0x40), alignment=PP_ALIGN.CENTER)

    out = os.path.join(os.path.dirname(__file__), "EDA-LIVESTOCK_営業資料_2026.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
