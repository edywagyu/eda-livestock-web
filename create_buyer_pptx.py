#!/usr/bin/env python3
"""EDA WAGYU — Buyer's Guide 2026 PPTX Generator. 10-slide buyer deck (revised)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FOREST = RGBColor(0x0F, 0x3D, 0x2E)
FOREST_DEEP = RGBColor(0x0A, 0x2D, 0x21)
GOLD = RGBColor(0xD4, 0xA9, 0x3B)
GOLD_SOFT = RGBColor(0xB8, 0x93, 0x2F)
CREAM = RGBColor(0xFA, 0xF7, 0xF0)
PAPER = RGBColor(0xF5, 0xEF, 0xE3)
BONE = RGBColor(0xEF, 0xE8, 0xD7)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_MUTED = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_82 = RGBColor(0xD1, 0xD1, 0xD1)
WHITE_55 = RGBColor(0x8C, 0x8C, 0x8C)

FONT_DISPLAY = "Arial Black"
FONT_BODY = "Arial"
FONT_JP = "Yu Gothic"

SLD_W = Inches(13.333)
SLD_H = Inches(7.5)

BASE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(BASE, "public", "images")


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_name=FONT_BODY,
             font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, spacing_before=0, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.space_before = Pt(spacing_before)
    p.line_spacing = line_spacing
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def add_rich_text(slide, left, top, width, height, segments, alignment=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.line_spacing = line_spacing
    for i, seg in enumerate(segments):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            if seg.get("newline"):
                p = tf.add_paragraph()
                p.alignment = alignment
                p.line_spacing = line_spacing
                if seg.get("spacing_before"):
                    p.space_before = Pt(seg["spacing_before"])
                run = p.add_run()
            else:
                run = p.add_run()
        run.text = seg.get("text", "")
        run.font.name = seg.get("font", FONT_BODY)
        run.font.size = Pt(seg.get("size", 14))
        run.font.color.rgb = seg.get("color", WHITE)
        run.font.bold = seg.get("bold", False)
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    return shape


def add_line_shape(slide, left, top, width, color=GOLD, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_footer(slide, num, dark=True):
    opacity_color = RGBColor(0x80, 0x80, 0x80) if dark else RGBColor(0x99, 0x99, 0x99)
    add_text(slide, Inches(0.6), Inches(7.05), Inches(3), Inches(0.35),
             "EDA WAGYU — BUYER'S GUIDE", FONT_DISPLAY, 9, opacity_color, alignment=PP_ALIGN.LEFT)
    add_text(slide, Inches(9.7), Inches(7.05), Inches(3), Inches(0.35),
             f"{num:02d}", FONT_DISPLAY, 9, opacity_color, alignment=PP_ALIGN.RIGHT)


def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image if file exists, otherwise add placeholder rect."""
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)
    else:
        w = width or Inches(1)
        h = height or Inches(1)
        add_rect(slide, left, top, w, h, fill_color=RGBColor(0x33, 0x33, 0x33), border_color=GOLD, border_width=Pt(1))


def main():
    prs = Presentation()
    prs.slide_width = SLD_W
    prs.slide_height = SLD_H
    blank = prs.slide_layouts[6]

    # ════════════════════════════════════════════
    # SLIDE 01 — COVER (Simple)
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST)

    # Logo
    logo_path = os.path.join(IMG, "logo", "eda-livestock-logo-transparent.png")
    add_image_safe(s, logo_path, Inches(5.9), Inches(1.8), width=Inches(1.5))

    # Title
    add_rich_text(s, Inches(2), Inches(3.5), Inches(9.3), Inches(2.0), [
        {"text": "Eda Wagyu", "font": FONT_DISPLAY, "size": 60, "color": WHITE, "bold": True},
        {"text": "Buyer's Guide 2026", "font": FONT_DISPLAY, "size": 60, "color": GOLD, "bold": True, "newline": True},
    ], alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_rich_text(s, Inches(3), Inches(5.5), Inches(7.3), Inches(1.0), [
        {"text": "Antibiotic-free, hormone-free Japanese Black Wagyu", "font": FONT_BODY, "size": 14, "color": WHITE_55},
        {"text": "Direct from Miyazaki, Japan", "font": FONT_BODY, "size": 14, "color": WHITE_55, "newline": True},
    ], alignment=PP_ALIGN.CENTER)

    # Confidential
    add_text(s, Inches(8), Inches(7.05), Inches(5), Inches(0.35),
             "CONFIDENTIAL — EDA LIVESTOCK CO., LTD.", FONT_DISPLAY, 8,
             RGBColor(0x40, 0x40, 0x40), alignment=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════════
    # SLIDE 02 — ABOUT US / EDA LIVESTOCK
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST_DEEP)

    # Left half — Eda working photo
    eda_photo = os.path.join(IMG, "farm", "cattle-08.jpg")
    if os.path.exists(eda_photo):
        s.shapes.add_picture(eda_photo, Inches(0), Inches(0),
                              width=Inches(6.8), height=SLD_H)

    # Right half — content
    add_text(s, Inches(7.2), Inches(0.9), Inches(5), Inches(0.3),
             "— ABOUT US —", FONT_DISPLAY, 10, GOLD)

    add_rich_text(s, Inches(7.2), Inches(1.3), Inches(5.5), Inches(1.4), [
        {"text": "Eda Livestock", "font": FONT_DISPLAY, "size": 36, "color": WHITE, "bold": True},
        {"text": "Co., Ltd.", "font": FONT_DISPLAY, "size": 22, "color": GOLD, "bold": True, "newline": True},
    ])

    add_rich_text(s, Inches(7.2), Inches(2.85), Inches(5.5), Inches(1.0), [
        {"text": "Founded in 2023 by ", "font": FONT_BODY, "size": 11, "color": WHITE_82},
        {"text": "Tomoki Eda", "font": FONT_BODY, "size": 11, "color": GOLD, "bold": True},
        {"text": ", the 4th-generation successor of Eda Livestock — a 70-year cattle-farming heritage in Miyazaki, Japan.", "font": FONT_BODY, "size": 11, "color": WHITE_82},
    ], line_spacing=1.7)

    add_rich_text(s, Inches(7.2), Inches(3.85), Inches(5.5), Inches(1.0), [
        {"text": "We raise Japanese Black Wagyu using ", "font": FONT_BODY, "size": 11, "color": WHITE_82},
        {"text": "circular feeding", "font": FONT_BODY, "size": 11, "color": GOLD, "bold": True},
        {"text": ", zero antibiotics, and complete farm-to-table traceability — pioneering the world's first organic-certified Wagyu beef.", "font": FONT_BODY, "size": 11, "color": WHITE_82},
    ], line_spacing=1.7)

    # Quote box
    add_rect(s, Inches(7.2), Inches(5.0), Pt(2), Inches(0.85), fill_color=GOLD)
    add_rect(s, Inches(7.25), Inches(5.0), Inches(5.45), Inches(0.85),
             fill_color=RGBColor(0x16, 0x42, 0x2E))
    add_text(s, Inches(7.4), Inches(5.05), Inches(5.2), Inches(0.75),
             "\"Not the rank, but how they're raised. We believe true quality comes from the feed, the care, and the methods behind the meat.\"",
             FONT_BODY, 10, WHITE, line_spacing=1.6)

    # CEO line
    add_line_shape(s, Inches(7.2), Inches(6.05), Inches(5.5), color=RGBColor(0x4A, 0x6B, 0x55))
    # CEO info
    add_text(s, Inches(7.2), Inches(6.2), Inches(3), Inches(0.3),
             "Tomoki Eda", FONT_DISPLAY, 14, WHITE, bold=True)
    add_text(s, Inches(7.2), Inches(6.55), Inches(4), Inches(0.25),
             "CEO · 4TH GENERATION REPRESENTATIVE", FONT_DISPLAY, 9, GOLD)
    # 2023 founded
    add_text(s, Inches(11.5), Inches(6.18), Inches(1.2), Inches(0.5),
             "2023", FONT_DISPLAY, 26, GOLD, alignment=PP_ALIGN.RIGHT)
    add_text(s, Inches(11.0), Inches(6.65), Inches(1.7), Inches(0.25),
             "FOUNDED", FONT_DISPLAY, 9, RGBColor(0x90, 0x90, 0x90), alignment=PP_ALIGN.RIGHT)

    add_footer(s, 2, dark=True)

    # ════════════════════════════════════════════
    # SLIDE 03 — WHY EDA WAGYU
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)

    # Left half
    add_text(s, Inches(0.7), Inches(1.2), Inches(4), Inches(0.3),
             "— WHY EDA WAGYU —", FONT_DISPLAY, 10, GOLD)

    add_rich_text(s, Inches(0.7), Inches(1.6), Inches(5), Inches(1.5), [
        {"text": "What Makes Us", "font": FONT_DISPLAY, "size": 38, "color": FOREST, "bold": True},
        {"text": "Different.", "font": FONT_DISPLAY, "size": 38, "color": GOLD, "bold": True, "newline": True},
    ])

    add_text(s, Inches(0.7), Inches(3.2), Inches(5.3), Inches(1.2),
             "Unlike conventional Wagyu brands defined by region (Kobe, Matsusaka, Omi), Eda Wagyu is defined by how the cattle are raised — our feed and farming practices.",
             FONT_BODY, 12, INK_MUTED, line_spacing=1.8)

    add_text(s, Inches(0.7), Inches(4.3), Inches(5.3), Inches(1.0),
             "This means your customers get a consistent, transparent product with a story they can verify — not just a geographic label.",
             FONT_BODY, 12, INK_MUTED, line_spacing=1.8)

    # 1% callout
    add_rect(s, Inches(0.7), Inches(5.6), Inches(5), Inches(0.5), fill_color=FOREST)
    add_text(s, Inches(0.9), Inches(5.65), Inches(4.6), Inches(0.4),
             "ONLY 1% OF JAPANESE FARMS ARE ADDITIVE-FREE", FONT_DISPLAY, 10, GOLD)

    # Right half — green panel
    add_rect(s, Inches(6.667), Inches(0), Inches(6.666), SLD_H, fill_color=FOREST)

    points = [
        ("Antibiotic & Hormone Free", "Zero antibiotics, growth promoters, or hormones throughout the entire lifecycle."),
        ("Circular Feed System", "Feed produced from our own compost cycle — not imported. Structural cost advantage."),
        ("Grade A3+ · 28 Months+", "Minimum A3 grade, 28+ months fattening. Grade and region specification available."),
        ("World's Only Organic Wagyu", "Premium EdaWagyu holds JAS Organic + EU Leaf equivalence — first and only in the world."),
        ("Direct Export · JPY Settlement", "No trading house intermediaries. Direct from farm, settled in Japanese Yen."),
    ]
    y_start = 0.9
    for i, (title, desc) in enumerate(points):
        y = y_start + i * 1.25
        add_line_shape(s, Inches(7.0), Inches(y), Pt(2), color=GOLD)
        add_rect(s, Inches(6.97), Inches(y), Pt(2), Inches(0.9), fill_color=GOLD)
        add_text(s, Inches(7.2), Inches(y + 0.05), Inches(5.5), Inches(0.3),
                 title, FONT_DISPLAY, 14, WHITE, bold=True)
        add_text(s, Inches(7.2), Inches(y + 0.35), Inches(5.5), Inches(0.6),
                 desc, FONT_BODY, 10, WHITE_82, line_spacing=1.6)

    add_footer(s, 3, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 04 — CUT LINEUP (text-based, no photos)
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.7), Inches(0.5), Inches(5), Inches(0.3),
             "— EXPORT CUT LINEUP —", FONT_DISPLAY, 10, GOLD)

    add_rich_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.6), [
        {"text": "Six Primal Cuts — ", "font": FONT_DISPLAY, "size": 26, "color": FOREST, "bold": True},
        {"text": "Chilled or Frozen, Vacuum-Packed", "font": FONT_DISPLAY, "size": 26, "color": GOLD, "bold": True},
    ])

    add_text(s, Inches(0.7), Inches(1.4), Inches(11), Inches(0.3),
             "All cuts available chilled (0–4°C) or frozen (-18°C). Vacuum-pack or skin-pack. Custom portioning on request.",
             FONT_BODY, 10, INK_MUTED)

    cuts = [
        ("01", "Striploin", "サーロイン", "Rich marbling in a fine web pattern. The king of steaks.", "BEST SELLER", True),
        ("02", "Ribeye", "リブロース / リブアイ", "Generous marbling with buttery flavor. Excellent for roasting.", "POPULAR", True),
        ("03", "Tenderloin / Fillet", "ヒレ / フィレ", "Only ~3kg per head. The most tender cut — perfect for chateaubriand.", "PREMIUM", True),
        ("04", "Round (Momo)", "モモ（ウチモモ · ランプ · イチボ）", "Lean with concentrated umami. Versatile — roast, steak, shabu-shabu.", "VERSATILE", False),
        ("05", "Tomahawk", "トマホーク（骨付きリブロース）", "Bone-in ribeye with long rib bone. Dramatic presentation cut.", "SHOWPIECE", True),
        ("06", "Skin Pack Portions", "スキンパック（個包装）", "Pre-portioned 150g–300g. Retail-ready for high-end supermarkets.", "RETAIL READY", True),
    ]

    row_w = Inches(6.0)
    row_h = Inches(1.45)
    for i, (num, name, jp, desc, badge, has_badge) in enumerate(cuts):
        col = i % 2
        row = i // 2
        x = Inches(0.7) + col * (row_w + Inches(0.3))
        y = Inches(2.0) + row * (row_h + Inches(0.15))

        # Card background (cream)
        add_rect(s, x, y, row_w, row_h, fill_color=CREAM)
        # Gold left border
        add_rect(s, x, y, Pt(2), row_h, fill_color=GOLD)
        # Number
        add_text(s, x + Inches(0.1), y + Inches(0.3), Inches(0.8), Inches(0.6),
                 num, FONT_DISPLAY, 28, GOLD, bold=True)
        # Name
        add_text(s, x + Inches(1.0), y + Inches(0.18), Inches(3.4), Inches(0.3),
                 name, FONT_DISPLAY, 16, FOREST, bold=True)
        # JP
        add_text(s, x + Inches(1.0), y + Inches(0.45), Inches(3.4), Inches(0.25),
                 jp, FONT_JP, 9, GOLD_SOFT)
        # Desc
        add_text(s, x + Inches(1.0), y + Inches(0.72), Inches(4.3), Inches(0.7),
                 desc, FONT_BODY, 9, INK_MUTED, line_spacing=1.5)
        # Badge
        if has_badge:
            badge_x = x + row_w - Inches(1.05)
            add_rect(s, badge_x, y + Inches(0.18), Inches(0.95), Inches(0.22), fill_color=FOREST)
            add_text(s, badge_x, y + Inches(0.19), Inches(0.95), Inches(0.22),
                     badge, FONT_DISPLAY, 8, GOLD, alignment=PP_ALIGN.CENTER)
        else:
            add_text(s, x + row_w - Inches(1.05), y + Inches(0.19), Inches(0.95), Inches(0.22),
                     badge, FONT_DISPLAY, 8, INK_MUTED, alignment=PP_ALIGN.CENTER)

    add_footer(s, 4, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 04 — SPEC & GRADING
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.7), Inches(0.8), Inches(5), Inches(0.3),
             "— PRODUCT SPECIFICATION —", FONT_DISPLAY, 10, GOLD)

    add_rich_text(s, Inches(0.7), Inches(1.1), Inches(10), Inches(0.6), [
        {"text": "Two Brands, ", "font": FONT_DISPLAY, "size": 26, "color": FOREST, "bold": True},
        {"text": "One Standard of Excellence", "font": FONT_DISPLAY, "size": 26, "color": GOLD, "bold": True},
    ])

    # Table header
    cols = [("Specification", 2.8, CREAM, INK), ("Eda Wagyu", 3.0, FOREST, WHITE),
            ("Premium EdaWagyu", 3.0, FOREST, WHITE), ("Conventional Wagyu", 2.8, CREAM, INK)]
    x_pos = 0.7
    for label, w, bg_c, txt_c in cols:
        add_rect(s, Inches(x_pos), Inches(1.8), Inches(w), Inches(0.4), fill_color=bg_c)
        add_text(s, Inches(x_pos + 0.1), Inches(1.82), Inches(w - 0.2), Inches(0.35),
                 label, FONT_DISPLAY, 10, txt_c, bold=True,
                 alignment=PP_ALIGN.CENTER if bg_c == FOREST else PP_ALIGN.LEFT)
        x_pos += w

    # Table rows
    rows = [
        ("Breed", "Japanese Black (Kuroge)", "Japanese Black (Kuroge)", "Japanese Black"),
        ("Antibiotics / Hormones", "✓ Zero", "✓ Zero", "Used"),
        ("Feed", "Circular (domestic)", "100% Organic-certified", "Imported compound"),
        ("Minimum Grade", "A3+", "A3+", "Varies"),
        ("Fattening Period", "28+ months", "31+ months", "24–26 months"),
        ("Regions", "Miyazaki · Kagoshima · Tohoku", "Miyazaki only (Eda Farm)", "Single prefecture"),
        ("Organic Certification", "—", "✓ JAS + EU Leaf + USDA", "—"),
        ("HALAL", "✓ Available", "Not available", "Limited"),
        ("Animal Welfare", "✓ Compliant", "✓ Organic standard", "No standard"),
        ("Traceability", "✓ Full (birth–table)", "✓ Full (birth–table)", "Partial"),
    ]
    for ri, (spec, eda, prem, conv) in enumerate(rows):
        y = Inches(2.25 + ri * 0.38)
        vals = [spec, eda, prem, conv]
        x_pos = 0.7
        for ci, (_, w, _, _) in enumerate(cols):
            color = INK if ci == 0 else (FOREST if "✓" in vals[ci] else (RGBColor(0xCC, 0xCC, 0xCC) if vals[ci] in ("Used", "—", "Limited", "No standard", "Not available") else INK))
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            fnt_bold = "✓" in vals[ci]
            add_text(s, Inches(x_pos + 0.1), y, Inches(w - 0.2), Inches(0.35),
                     vals[ci], FONT_BODY, 10, color, bold=fnt_bold, alignment=align)
            # Row divider
            add_line_shape(s, Inches(x_pos), y + Inches(0.35), Inches(w), color=RGBColor(0xEE, 0xEE, 0xEE))
            x_pos += w

    # Note bar
    add_rect(s, Inches(0.7), Inches(6.2), Inches(7), Inches(0.4), fill_color=FOREST)
    add_text(s, Inches(0.9), Inches(6.22), Inches(6.6), Inches(0.35),
             "PREMIUM EDAWAGYU IS AVAILABLE EXCLUSIVELY TO EXISTING EDA WAGYU CUSTOMERS",
             FONT_DISPLAY, 10, GOLD)

    add_footer(s, 5, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 05 — LOGISTICS & PHASES
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST)

    add_text(s, Inches(0.7), Inches(0.6), Inches(5), Inches(0.3),
             "— LOGISTICS & SHIPPING —", FONT_DISPLAY, 10, GOLD)

    add_rich_text(s, Inches(0.7), Inches(0.9), Inches(10), Inches(0.6), [
        {"text": "Farm to Your Designated Port in ", "font": FONT_DISPLAY, "size": 26, "color": WHITE, "bold": True},
        {"text": "~1 Month", "font": FONT_DISPLAY, "size": 26, "color": GOLD, "bold": True},
    ])

    add_text(s, Inches(0.7), Inches(1.5), Inches(10), Inches(0.3),
             "Chilled or frozen — you choose. We handle export documents, customs clearance, and freight to your designated port.",
             FONT_BODY, 10, WHITE_55)

    # 3 Phase cards
    phases = [
        ("PHASE 1 — ORDER", "Order Confirmation",
         ["Purchase order confirmed", "Invoice issued (JPY)", "30% deposit received", "Slaughter scheduled"]),
        ("PHASE 2 — PREPARATION", "Documents & Shipping",
         ["Health Certificate issued", "Packing List prepared", "Commercial Invoice finalized", "HALAL cert (if applicable)", "70% balance payment received"]),
        ("PHASE 3 — DELIVERY", "Shipment & Arrival",
         ["AWB (Air Waybill) issued", "EQC (Export Quarantine Cert)", "Certificate of Origin", "Freight to designated port", "Arrival confirmation"]),
    ]
    card_w = Inches(3.8)
    for i, (label, title, items) in enumerate(phases):
        x = Inches(0.7) + i * (card_w + Inches(0.25))
        # Card bg
        add_rect(s, x, Inches(2.0), card_w, Inches(2.8),
                 fill_color=RGBColor(0x12, 0x45, 0x35), border_color=RGBColor(0x1A, 0x55, 0x40), border_width=Pt(1))
        # Gold top border
        add_rect(s, x, Inches(2.0), card_w, Pt(3), fill_color=GOLD)
        # Label
        add_text(s, x + Inches(0.2), Inches(2.15), card_w - Inches(0.4), Inches(0.25),
                 label, FONT_DISPLAY, 10, GOLD)
        # Title
        add_text(s, x + Inches(0.2), Inches(2.4), card_w - Inches(0.4), Inches(0.3),
                 title, FONT_DISPLAY, 14, WHITE, bold=True)
        # Items
        for j, item in enumerate(items):
            add_text(s, x + Inches(0.2), Inches(2.8 + j * 0.3), card_w - Inches(0.4), Inches(0.25),
                     f"→ {item}", FONT_BODY, 9, WHITE_82)

    # Bottom details
    details = [
        ("TEMPERATURE", "Chilled: 0–4°C or Frozen: -18°C.\nBuyer specifies at order."),
        ("PACKAGING", "Vacuum-pack (standard) or skin-pack (retail).\nStyrofoam box with gel ice / dry ice."),
        ("SHELF LIFE", "110 days from production.\nApprox. 90 days remaining at arrival."),
    ]
    for i, (title, desc) in enumerate(details):
        x = Inches(0.7) + i * Inches(4.05)
        add_rect(s, x, Inches(5.1), Pt(2), Inches(0.8), fill_color=GOLD)
        add_text(s, x + Inches(0.2), Inches(5.15), Inches(3.6), Inches(0.2),
                 title, FONT_DISPLAY, 11, GOLD)
        add_text(s, x + Inches(0.2), Inches(5.4), Inches(3.6), Inches(0.6),
                 desc, FONT_BODY, 10, WHITE_82, line_spacing=1.6)

    add_footer(s, 6, dark=True)

    # ════════════════════════════════════════════
    # SLIDE 06 — COMMERCIAL TERMS
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)

    add_text(s, Inches(0.7), Inches(0.5), Inches(5), Inches(0.3),
             "— COMMERCIAL TERMS —", FONT_DISPLAY, 10, GOLD)

    add_rich_text(s, Inches(0.7), Inches(0.8), Inches(10), Inches(0.6), [
        {"text": "Transparent Terms, ", "font": FONT_DISPLAY, "size": 26, "color": FOREST, "bold": True},
        {"text": "Built for Long-Term Partnership", "font": FONT_DISPLAY, "size": 26, "color": GOLD, "bold": True},
    ])

    # Left column
    add_text(s, Inches(0.7), Inches(1.7), Inches(5), Inches(0.35),
             "Order & Payment", FONT_DISPLAY, 18, FOREST, bold=True)

    left_terms = [
        ("CURRENCY", "Japanese Yen (JPY) — all invoices in ¥"),
        ("PAYMENT", "Deposit 30% at order confirmation\nBalance 70% before shipping\nThis rule applies to all clients without exception"),
        ("LEAD TIME", "Approx. 4 weeks from deposit to port delivery"),
        ("MOQ", "200 kg per order"),
        ("PRICING", "Based on JPY/kg · Grade and region affect price"),
        ("EXCLUSIVITY", "Currently exclusive in 3 countries.\nInquire separately for distribution rights."),
    ]
    y = 2.2
    for label, value in left_terms:
        add_text(s, Inches(0.7), Inches(y), Inches(1.4), Inches(0.25),
                 label, FONT_DISPLAY, 9, GOLD_SOFT)
        add_text(s, Inches(2.2), Inches(y), Inches(4.2), Inches(0.65),
                 value, FONT_BODY, 11, INK, line_spacing=1.5)
        add_line_shape(s, Inches(0.7), Inches(y + 0.55), Inches(5.7), color=RGBColor(0xE0, 0xE0, 0xE0))
        y += 0.65

    # Right column
    add_text(s, Inches(6.8), Inches(1.7), Inches(5), Inches(0.35),
             "Shipping & Delivery", FONT_DISPLAY, 18, FOREST, bold=True)

    right_terms = [
        ("INCOTERMS", "FOB or CIF — buyer's choice"),
        ("FREIGHT", "Air freight, chilled or frozen"),
        ("DELIVERY", "Approx. 4 weeks to designated port\nCustoms clearance by buyer"),
        ("SHELF LIFE", "110 days from production\n~90 days remaining at arrival"),
        ("CUSTOM CUT", "Thickness, portioning, and trim to spec."),
        ("GRADE SELECT", "A3/A4/A5 specification available.\nRegion lock = price premium."),
    ]
    y = 2.2
    for label, value in right_terms:
        add_text(s, Inches(6.8), Inches(y), Inches(1.5), Inches(0.25),
                 label, FONT_DISPLAY, 9, GOLD_SOFT)
        add_text(s, Inches(8.4), Inches(y), Inches(4.5), Inches(0.55),
                 value, FONT_BODY, 11, INK, line_spacing=1.5)
        add_line_shape(s, Inches(6.8), Inches(y + 0.5), Inches(6.1), color=RGBColor(0xE0, 0xE0, 0xE0))
        y += 0.58

    add_footer(s, 7, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 07 — CERTIFICATIONS
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST_DEEP)

    add_text(s, Inches(0.7), Inches(0.8), Inches(5), Inches(0.3),
             "— CERTIFICATIONS —", FONT_DISPLAY, 10, GOLD)

    add_rich_text(s, Inches(0.7), Inches(1.1), Inches(10), Inches(0.6), [
        {"text": "Trust, Verified by ", "font": FONT_DISPLAY, "size": 26, "color": WHITE, "bold": True},
        {"text": "International Standards", "font": FONT_DISPLAY, "size": 26, "color": GOLD, "bold": True},
    ])

    certs = [
        ("Eda Wagyu", "BOTH BRANDS", "Proprietary brand cert.\nAntibiotic-free, hormone-free.", True, "goods/logo-sticker-transparent.png"),
        ("HALAL", "EDA WAGYU ONLY", "Certified for Malaysia, Indonesia,\nand Islamic markets.", False, None),
        ("USDA Organic", "PREMIUM ONLY", "Application in progress.\nTargeting US organic market.", False, None),
        ("EU Leaf", "PREMIUM ONLY", "Approved. Organic-labeled\nexport to all EU states.", False, None),
        ("JAS Organic", "PREMIUM ONLY", "World's first organic-certified\nKuroge Wagyu. Jan 2026.", True, None),
    ]

    card_w = Inches(2.2)
    for i, (name, tag, desc, gold_bdr, img) in enumerate(certs):
        x = Inches(0.7) + i * (card_w + Inches(0.2))
        border_c = GOLD if gold_bdr else RGBColor(0x1A, 0x55, 0x40)
        bg_c = RGBColor(0x14, 0x3A, 0x2A) if gold_bdr else RGBColor(0x0D, 0x30, 0x22)
        add_rect(s, x, Inches(1.8), card_w, Inches(2.2), fill_color=bg_c, border_color=border_c, border_width=Pt(1))
        # Logo for first card
        if img:
            img_path = os.path.join(IMG, img)
            if os.path.exists(img_path):
                add_image_safe(s, img_path, x + Inches(0.7), Inches(1.95), width=Inches(0.8))
        # Card title
        add_text(s, x, Inches(2.7), card_w, Inches(0.3), name,
                 FONT_DISPLAY, 13, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Tag
        tag_y = Inches(3.0)
        add_text(s, x, tag_y, card_w, Inches(0.2), tag,
                 FONT_DISPLAY, 8, GOLD, alignment=PP_ALIGN.CENTER)
        # Desc
        add_text(s, x + Inches(0.15), Inches(3.25), card_w - Inches(0.3), Inches(0.7), desc,
                 FONT_BODY, 9, WHITE_82, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    # Bottom callout
    add_rect(s, Inches(0.7), Inches(4.2), Inches(11.9), Inches(1.6),
             border_color=GOLD, border_width=Pt(1))

    callouts = [
        ("For Your Menu", "\"Antibiotic-free,\" \"Hormone-free,\" \"Organic-certified\" — menu claims your guests actively seek."),
        ("For Your Brand", "Sustainability, traceability, and animal welfare are expected by luxury consumers."),
        ("For Your Compliance", "EU import regulations, HALAL requirements, organic labeling — we handle the paperwork."),
    ]
    for i, (title, desc) in enumerate(callouts):
        x = Inches(0.9) + i * Inches(3.9)
        add_text(s, x, Inches(4.35), Inches(3.5), Inches(0.25),
                 title.upper(), FONT_DISPLAY, 10, GOLD)
        add_text(s, x, Inches(4.6), Inches(3.5), Inches(0.9),
                 desc, FONT_BODY, 10, RGBColor(0xB0, 0xB0, 0xB0), line_spacing=1.6)

    add_footer(s, 8, dark=True)

    # ════════════════════════════════════════════
    # SLIDE 08 — TRUST / WHERE SERVED
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)

    add_text(s, Inches(0.7), Inches(0.6), Inches(5), Inches(0.3),
             "— WHERE EDA WAGYU IS SERVED —", FONT_DISPLAY, 10, GOLD)

    add_rich_text(s, Inches(0.7), Inches(0.9), Inches(10), Inches(0.6), [
        {"text": "Trusted by ", "font": FONT_DISPLAY, "size": 26, "color": FOREST, "bold": True},
        {"text": "The World's Finest Hotels", "font": FONT_DISPLAY, "size": 26, "color": GOLD, "bold": True},
    ])

    hotels = [
        ("THE RITZ-CARLTON", "Tokyo · Fukuoka · Hong Kong · New York", True),
        ("FOUR SEASONS", "Tokyo at Otemachi", True),
        ("THE PENINSULA", "Tokyo", True),
        ("MANDARIN ORIENTAL", "Tokyo", True),
        ("TOKYO STATION HOTEL", "National Heritage · Tokyo", False),
        ("EDITION GINZA", "Marriott Luxury · Ginza", False),
        ("CITY'SUPER", "All stores · Hong Kong", False),
        ("MITSUKOSHI · KEIO", "Ginza · Nihonbashi · Shinjuku", False),
    ]
    card_w = Inches(2.8)
    for i, (name, loc, feat) in enumerate(hotels):
        col = i % 4
        row = i // 4
        x = Inches(0.7) + col * (card_w + Inches(0.2))
        y = Inches(1.7) + row * Inches(1.3)
        border_c = GOLD if feat else BONE
        bg_c = RGBColor(0xF8, 0xF3, 0xE8) if feat else None
        add_rect(s, x, y, card_w, Inches(1.1), fill_color=bg_c, border_color=border_c, border_width=Pt(1))
        add_text(s, x, y + Inches(0.25), card_w, Inches(0.3), name,
                 FONT_BODY, 13, FOREST, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.6), card_w, Inches(0.25), loc,
                 FONT_BODY, 8, INK_MUTED, alignment=PP_ALIGN.CENTER)

    # Stats
    stats = [("21", "LOCATIONS"), ("13", "COUNTRIES"), ("4", "CONTINENTS")]
    for i, (num, label) in enumerate(stats):
        x = Inches(0.7) + i * Inches(1.5)
        add_text(s, x, Inches(4.6), Inches(1.3), Inches(0.7), num,
                 FONT_DISPLAY, 42, GOLD)
        add_text(s, x, Inches(5.2), Inches(1.3), Inches(0.25), label,
                 FONT_DISPLAY, 9, INK_MUTED)

    add_text(s, Inches(6.5), Inches(4.8), Inches(6), Inches(0.8),
             "Also served at ANA International First Class, Tokyo DisneySea, WITH THE STYLE Fukuoka, and fine-dining restaurants across Europe and Southeast Asia.",
             FONT_BODY, 11, INK_MUTED, alignment=PP_ALIGN.RIGHT, line_spacing=1.7)

    add_footer(s, 9, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 10 — BRANDED GOODS (3x2 Editorial Grid)
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)

    add_text(s, Inches(0.7), Inches(0.4), Inches(8), Inches(0.3),
             "— BRANDED MATERIALS · SOLD SEPARATELY —", FONT_DISPLAY, 10, GOLD)

    add_rich_text(s, Inches(0.7), Inches(0.75), Inches(12), Inches(0.6), [
        {"text": "Elevate Your Presentation with ", "font": FONT_DISPLAY, "size": 24, "color": FOREST, "bold": True},
        {"text": "Eda Wagyu Originals", "font": FONT_DISPLAY, "size": 24, "color": GOLD, "bold": True},
    ])

    add_text(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.3),
             "Optional brand materials available for purchase. Strengthen the in-store and dining experience with authentic, hand-crafted Eda Wagyu items.",
             FONT_BODY, 10, INK_MUTED)

    # 3x2 Grid of equal-sized cards
    items_3x2 = [
        # Row 1 — (name, jp, price, desc, img, tag, kind, fit_mode, bg_color)
        ("Original Branded Box", "江田和牛 オリジナル 化粧箱", "¥1,100 / BOX",
         "Premium green & gold-foil packaging with Hokusai-inspired wave and Mt. Fuji motif.",
         "goods/eda-box-green.jpg", "SIGNATURE", "cover", FOREST_DEEP),
        ("Brand Pamphlet", "3つ折りパンフレット（表・裏）", "¥440 / PIECE",
         "Bilingual 3-fold brochure (2 sides) detailing philosophy, farming, and certifications.",
         "goods/pamphlet-dual.jpg", "EN / JP", "contain", FOREST_DEEP),
        ("Eda Wagyu Certificate", "江田和牛 証明書", "¥550 / SHEET",
         "Individual carcass certificate (Cattle ID, Grade, Farm, Origin).",
         "certs/cert-clean.jpg", "AUTHENTIC", "contain", FOREST_DEEP),
        # Row 2
        ("Brand Sticker Set", "江田和牛 ロゴシール", "¥980 / SET",
         "Forest-green and gold seal stickers for packaging, gift wrap, and retail.",
         "goods/logo-sticker.jpg", "SET OF 24", "contain", FOREST_DEEP),
        ("Display Plate (Tatami)", "立て札 — 畳バージョン", "¥38,000 / PLATE",
         "Traditional tatami-textured display with embroidered \"江田和牛\" gold branding.",
         "goods/tatami-plate.jpg", "HANDMADE", "contain", PAPER),
        ("Display Plate (Wood)", "立て札 — 木製バージョン", "¥77,000 / PLATE",
         "Premium solid wood with engraved \"江田和牛\" — luxury statement piece.",
         "goods/wood-plate.jpg", "PREMIUM", "contain", PAPER),
    ]

    grid_left = Inches(0.6)
    grid_top = Inches(1.85)
    card_w = Inches(4.05)
    card_h = Inches(2.5)
    img_h = Inches(1.45)
    gap_x = Inches(0.1)
    gap_y = Inches(0.1)

    for i, (name, jp, price, desc, img_file, tag, fit_mode, bg_color) in enumerate(items_3x2):
        col = i % 3
        row = i // 3
        x = grid_left + col * (card_w + gap_x)
        y = grid_top + row * (card_h + gap_y)

        # Card border
        add_rect(s, x, y, card_w, card_h, fill_color=WHITE, border_color=GOLD, border_width=Pt(1))

        # Image area background
        add_rect(s, x, y, card_w, img_h, fill_color=bg_color)

        # Add image with proper fit
        if img_file:
            img_path = os.path.join(IMG, img_file)
            if os.path.exists(img_path):
                try:
                    from PIL import Image as PILImage
                    pil = PILImage.open(img_path)
                    ratio = pil.width / pil.height
                    target_w_in = float(card_w) / 914400
                    target_h_in = float(img_h) / 914400
                    container_ratio = target_w_in / target_h_in

                    if fit_mode == "contain":
                        # Show full image, fit inside container
                        if ratio > container_ratio:
                            # Image wider — fit by width
                            pic_w = target_w_in
                            pic_h = pic_w / ratio
                        else:
                            # Image taller — fit by height
                            pic_h = target_h_in
                            pic_w = pic_h * ratio
                    else:  # cover
                        if ratio > container_ratio:
                            pic_h = target_h_in
                            pic_w = pic_h * ratio
                        else:
                            pic_w = target_w_in
                            pic_h = pic_w / ratio

                    # Center image in container
                    cx = x + Inches((target_w_in - pic_w) / 2)
                    cy = y + Inches((target_h_in - pic_h) / 2)

                    if fit_mode == "cover":
                        # Crop using a clip box - but pptx doesn't support clipping easily
                        # Just place the image, may overflow slightly
                        s.shapes.add_picture(img_path, cx, cy, width=Inches(pic_w), height=Inches(pic_h))
                        # Cover up overflow with bg-colored rects on the sides
                        if pic_w > target_w_in:
                            overflow = (pic_w - target_w_in) / 2
                            add_rect(s, x - Inches(overflow), y, Inches(overflow), card_h, fill_color=CREAM)
                            add_rect(s, x + card_w, y, Inches(overflow), card_h, fill_color=CREAM)
                        if pic_h > target_h_in:
                            overflow = (pic_h - target_h_in) / 2
                            add_rect(s, x, y - Inches(overflow), card_w, Inches(overflow), fill_color=CREAM)
                            add_rect(s, x, y + img_h, card_w, Inches(overflow), fill_color=CREAM)
                    else:
                        # Contain: just place centered, full image visible
                        s.shapes.add_picture(img_path, cx, cy, width=Inches(pic_w), height=Inches(pic_h))
                except Exception:
                    s.shapes.add_picture(img_path, x, y, width=card_w, height=img_h)

        # Tag
        add_rect(s, x + Inches(0.1), y + Inches(0.08), Inches(0.7), Inches(0.2), fill_color=FOREST)
        add_text(s, x + Inches(0.1), y + Inches(0.09), Inches(0.7), Inches(0.2),
                 tag, FONT_DISPLAY, 7, GOLD, alignment=PP_ALIGN.CENTER)

        # Body text
        text_top = y + img_h + Inches(0.08)
        add_text(s, x + Inches(0.12), text_top, card_w - Inches(0.24), Inches(0.22),
                 name, FONT_DISPLAY, 11, FOREST, bold=True)
        add_text(s, x + Inches(0.12), text_top + Inches(0.22), card_w - Inches(0.24), Inches(0.18),
                 jp, FONT_JP, 8, GOLD_SOFT)
        add_text(s, x + Inches(0.12), text_top + Inches(0.42), card_w - Inches(0.24), Inches(0.22),
                 price, FONT_DISPLAY, 10, GOLD, bold=True)
        add_text(s, x + Inches(0.12), text_top + Inches(0.65), card_w - Inches(0.24), Inches(0.4),
                 desc, FONT_BODY, 8, INK_MUTED, line_spacing=1.4)
    add_footer(s, 10, dark=False)

    # ════════════════════════════════════════════
    # SLIDE 10 — HOW TO ORDER
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST_DEEP)

    add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.3),
             "— HOW TO ORDER —", FONT_DISPLAY, 10, GOLD, alignment=PP_ALIGN.CENTER)

    add_rich_text(s, Inches(1), Inches(1.2), Inches(11.3), Inches(0.6), [
        {"text": "From Inquiry to Delivery — ", "font": FONT_DISPLAY, "size": 26, "color": WHITE, "bold": True},
        {"text": "Simple & Direct", "font": FONT_DISPLAY, "size": 26, "color": GOLD, "bold": True},
    ], alignment=PP_ALIGN.CENTER)

    steps = [
        ("01", "Contact Us", "Email your interest, preferred cuts, volume, and destination. We reply within 24 hours."),
        ("02", "Quotation", "We send a JPY price list tailored to your order. Custom cuts and grades available."),
        ("03", "Deposit 30%", "Sign purchase order and send 30% deposit. Slaughter is scheduled upon receipt."),
        ("04", "Balance & Ship", "Pay remaining 70% before shipping. Documents + freight arranged to your designated port."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.8) + i * Inches(3.1)
        # Number box
        add_rect(s, x + Inches(0.85), Inches(2.1), Inches(0.75), Inches(0.75),
                 border_color=GOLD, border_width=Pt(2))
        add_text(s, x + Inches(0.85), Inches(2.15), Inches(0.75), Inches(0.7), num,
                 FONT_DISPLAY, 22, GOLD, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(s, x, Inches(3.1), Inches(2.5), Inches(0.3), title,
                 FONT_DISPLAY, 14, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Desc
        add_text(s, x, Inches(3.5), Inches(2.5), Inches(0.7), desc,
                 FONT_BODY, 10, WHITE_82, alignment=PP_ALIGN.CENTER, line_spacing=1.6)

    # Bottom callouts
    callouts = [
        ("First Order?", "We accommodate smaller first orders so you can test with your kitchen team before scaling."),
        ("Regular Buyer?", "Established accounts receive priority scheduling and access to Premium EdaWagyu (organic)."),
        ("Distributor?", "We partner with local importers and distributors. Exclusive territory rights negotiable."),
    ]
    for i, (title, desc) in enumerate(callouts):
        x = Inches(0.7) + i * Inches(4.1)
        add_rect(s, x, Inches(4.6), Inches(3.7), Pt(2), fill_color=GOLD)
        add_text(s, x, Inches(4.75), Inches(3.7), Inches(0.25), title.upper(),
                 FONT_DISPLAY, 10, GOLD)
        add_text(s, x, Inches(5.0), Inches(3.7), Inches(0.7), desc,
                 FONT_BODY, 10, WHITE_82, line_spacing=1.6)

    add_footer(s, 11, dark=True)

    # ════════════════════════════════════════════
    # SLIDE 11 — CONTACT
    # ════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, FOREST)

    # Logo
    add_image_safe(s, logo_path, Inches(5.9), Inches(1.5), width=Inches(1.5))

    add_rich_text(s, Inches(2), Inches(3.0), Inches(9.3), Inches(1.0), [
        {"text": "Ready to ", "font": FONT_DISPLAY, "size": 40, "color": WHITE, "bold": True},
        {"text": "Taste the Difference?", "font": FONT_DISPLAY, "size": 40, "color": GOLD, "bold": True},
    ], alignment=PP_ALIGN.CENTER)

    add_text(s, Inches(3.5), Inches(3.9), Inches(6.3), Inches(0.4),
             "Request a price list or schedule a call with our export team.",
             FONT_BODY, 13, WHITE_55, alignment=PP_ALIGN.CENTER)

    # Left contact
    add_line_shape(s, Inches(3.5), Inches(4.7), Inches(2.8), color=GOLD)
    add_text(s, Inches(3.5), Inches(4.85), Inches(2.8), Inches(0.25),
             "EXPORT INQUIRIES", FONT_DISPLAY, 11, GOLD)
    add_rich_text(s, Inches(3.5), Inches(5.15), Inches(2.8), Inches(1.2), [
        {"text": "Tomoki Eda — CEO / 4th Generation Representative", "font": FONT_BODY, "size": 12, "color": WHITE_82},
        {"text": "", "font": FONT_BODY, "size": 12, "color": WHITE_82, "newline": True},
        {"text": "backoffice@eda-livestock.com", "font": FONT_BODY, "size": 12, "color": GOLD, "newline": True},
        {"text": "", "font": FONT_BODY, "size": 12, "color": WHITE_82, "newline": True},
        {"text": "We respond within 24 hours.", "font": FONT_BODY, "size": 12, "color": WHITE_82, "newline": True},
    ], line_spacing=1.7)

    # Right contact
    add_line_shape(s, Inches(7.5), Inches(4.7), Inches(2.8), color=GOLD)
    add_text(s, Inches(7.5), Inches(4.85), Inches(2.8), Inches(0.25),
             "HEADQUARTERS", FONT_DISPLAY, 11, GOLD)
    add_rich_text(s, Inches(7.5), Inches(5.15), Inches(3.5), Inches(1.5), [
        {"text": "Eda Livestock Co., Ltd.", "font": FONT_BODY, "size": 12, "color": WHITE_82},
        {"text": "2025-2 Hirohara, Takaharu-cho", "font": FONT_BODY, "size": 12, "color": WHITE_82, "newline": True},
        {"text": "Miyazaki 889-4411, Japan", "font": FONT_BODY, "size": 12, "color": WHITE_82, "newline": True},
        {"text": "", "font": FONT_BODY, "size": 12, "color": WHITE_82, "newline": True},
        {"text": "backoffice@eda-livestock.com", "font": FONT_BODY, "size": 12, "color": GOLD, "newline": True},
        {"text": "eda-livestock.com · @eda_livestock", "font": FONT_BODY, "size": 12, "color": GOLD, "newline": True},
    ], line_spacing=1.7)

    # Copyright
    add_text(s, Inches(3), Inches(7.0), Inches(7.3), Inches(0.35),
             "© 2026 EDA-LIVESTOCK CO., LTD. — CONFIDENTIAL",
             FONT_DISPLAY, 8, RGBColor(0x40, 0x40, 0x40), alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # SAVE
    # ════════════════════════════════════════════
    out = os.path.join(BASE, "EDA-WAGYU_Buyers_Guide_2026.pptx")
    prs.save(out)
    print(f"Saved: {out} ({os.path.getsize(out):,} bytes)")


if __name__ == "__main__":
    main()
